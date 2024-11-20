from flask import Flask, render_template, request, redirect, url_for, session, jsonify
from flask_cors import CORS
from transformers import pipeline
import os
import textract
import requests
from ocr_processing import process_image
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from workers import txt2questions
from pptx import Presentation

app = Flask(__name__)
CORS(app)

app.secret_key = "Fg4bCUP3odF0ZvMgIS3wqJudc30Us0nv"

bart_summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
bert_summarizer = pipeline("summarization", model="bert-base-uncased")


@app.route('/')
def main():
    return render_template("index.html")


@app.route('/upload-image', methods=['POST'])
def upload_image():
    # Retrieve the uploaded file
    if 'image' in request.files and request.files['image'].filename != '':
        f = request.files['image']
        filename = f.filename
        f.save(filename)

        if filename.lower().endswith(('.jpg', '.jpeg', '.png')):

            # Retrieve radio button selections
            performance = request.form.get('performance')  # Faster or Better
            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')

            output_text = process_image(filename, performance)
            output_filename = f"{filename}_output.txt"

            with open(output_filename, 'w') as output_file:
                output_file.write(output_text)

            with open(output_filename, 'r', encoding="utf-8") as file:
                text = file.read()

            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")

        else:
            return jsonify(success=False, message="File uploaded successfully but it is not an image.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-pdf', methods=['POST'])
def upload_pdf():
    if 'pdf' in request.files and request.files['pdf'].filename != '':
        pdf_file = request.files['pdf']
        filename = pdf_file.filename
        pdf_file.save(filename)
        
        if filename.lower().endswith(('pdf')):
            extracted_text = textract.process(filename).decode('utf-8')
            text_filename = f"{filename}_text.txt"

            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)
    
            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')

            # with open(text_filename, 'w') as output_file:
            #     output_file.write(text)

            # with open(text_filename, 'r', encoding="utf-8") as file:
            #     text = file.read()

            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a PDF.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-ppt', methods=['POST'])
def upload_ppt():
    if 'ppt' in request.files and request.files['ppt'].filename != '':
        ppt_file = request.files['ppt']
        filename = ppt_file.filename
        ppt_file.save(filename)

        if filename.lower().endswith(('pptx')):

            def extract_text_from_ppt(filename):
                prs = Presentation(filename)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                text.append(paragraph.text)
                return "\n".join(text)

            extracted_text = extract_text_from_ppt(filename)
            text_filename = f"{filename}_text.txt"

            print(text_filename, "\n\n\n")

            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)
    

            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')
            
            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a PPT.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/upload-doc', methods=['POST'])
def upload_doc():
    if 'doc' in request.files and request.files['doc'].filename != '':
        doc_file = request.files['doc']
        filename = doc_file.filename
        doc_file.save(filename)
        
        if filename.lower().endswith(('doc', 'docx')):
            extracted_text = textract.process(filename).decode('utf-8')
            text_filename = f"{filename}_text.txt"
            
            with open(text_filename, "w", encoding="utf-8") as text_file:
                text_file.write(extracted_text)

            summary_type = request.form.get('summaryType')
            summary_length = request.form.get('summaryLength')
            
            max_input_length = 512  # Model's input token limit

            if summary_length == 'short':
                max_output_length = 50
            elif summary_length == 'medium':
                max_output_length = 130
            elif summary_length == 'long':
                max_output_length = 250
            else:
                max_output_length = 130

            truncated_text = extracted_text[:max_input_length]

            try:
                if summary_type == 'abstractive':
                    summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                elif summary_type == 'extractive':
                    summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
                else:
                    summary = "Invalid summary type selected."

                summary_filename = f"{filename}_summary.txt"
                
                with open(summary_filename, 'w', encoding="utf-8") as summary_file:
                    summary_file.write(summary)

                message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
                return jsonify(success=True, message=message, summary=summary)  # Send summary to frontend
            except Exception as e:
                return jsonify(success=False, message=f"Error during summarization: {e}")
            
        else:
            return jsonify(success=False, message="File uploaded successfully but it is not a DOC.")
        
    return jsonify(success=False, message="No file was uploaded.")



@app.route('/process_image', methods=['POST'])
def process_image_route():
    filename = request.form.get('filename')
    ocr_option = request.form.get('ocr_option')
    
    output_text = process_image(filename, ocr_option)
    output_filename = f"{filename}_output.txt"
    with open(output_filename, 'w') as output_file:
        output_file.write(output_text)
    
    return render_template("select_summary.html", filename=output_filename)


@app.route('/success', methods=['POST'])
def success():
    if request.method == 'POST':
        option = request.form.get('option')

        # if option == 'quiz':
        #     return render_template("select_quiz.html", filename=request.form.get('filename'))
        
        if option == 'image' and 'image' in request.files and request.files['image'].filename != '':
            f = request.files['image']
            filename = f.filename
            f.save(filename)
            
            if filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                return render_template("select_ocr.html", filename=filename)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not an image.")
        
        # Handle URL link submission
        elif option == 'link' and request.form.get('link'):
            link = request.form['link']
            try:
                # Fetch HTML content of the URL
                response = requests.get(link)
                response.raise_for_status()
                
                # Save HTML content to a file
                html_filename = "downloaded_link.html"
                with open(html_filename, "w", encoding="utf-8") as html_file:
                    html_file.write(response.text)
                
                # Process the saved HTML file with Textract
                extracted_text = textract.process(html_filename, method='html').decode('utf-8')
                text_filename = "link_extracted_text.txt"
                with open(text_filename, "w", encoding="utf-8") as text_file:
                    text_file.write(extracted_text)
                return render_template("select_summary.html", filename=text_filename)
            except requests.exceptions.RequestException as e:
                message = f"Error fetching the URL: {e}"
            except textract.exceptions.ShellError as e:
                message = f"Error processing the HTML file with Textract: {e}"
            except Exception as e:
                message = f"An unexpected error occurred: {e}"
            return render_template("ack.html", message=message)

        # Handle PDF file upload
        elif option == 'pdf' and 'pdf' in request.files and request.files['pdf'].filename != '':
            pdf_file = request.files['pdf']
            filename = pdf_file.filename
            pdf_file.save(filename)
            if filename.lower().endswith(('pdf')):
                try:
                    extracted_text = textract.process(filename).decode('utf-8')
                    text_filename = f"{filename}_text.txt"
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the PDF file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a pdf.")
        
        # Handle DOC file upload
        elif option == 'doc' and 'doc' in request.files and request.files['doc'].filename != '':
            doc_file = request.files['doc']
            filename = doc_file.filename
            doc_file.save(filename)
            if filename.lower().endswith(('doc', 'docx')):
                try:
                    extracted_text = textract.process(filename).decode('utf-8')
                    text_filename = f"{filename}_text.txt"
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the DOC file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a doc.")

        
        elif option == 'quiz' and 'quiz' in request.files and request.files['quiz'].filename != '':
            UPLOAD_STATUS = False
            global questions
            questions = dict()

            quiz_file = request.files['quiz']
            # # print("\n\n\n\n\n")
            # # print("Quiz File?:", quiz_file)
            # # print("\n\n\n\n\n")
            filename = quiz_file.filename
            # # print("\n\n\n\n\n")
            # # print("Filename:", filename)
            # # print("\n\n\n\n\n")
            quiz_file.save(filename)

            if filename.lower().endswith(('txt')):
                try:
                    with open(filename, 'r') as file:
                        text = file.read()

                    questions = txt2questions(text)
                    
                    # print("\n\n\n\n\n")
                    # print(questions)
                    # print("\n\n\n\n\n")

                # File upload + convert success
                    if text is not None:
                        UPLOAD_STATUS = True

                    # print("\n\n\n\n\n")
                    # print("Reached Debug Point 1")
                    # print("\n\n\n\n\n")

                    # return redirect("/quiz")
                    return render_template('quiz.html', uploaded=UPLOAD_STATUS, questions=questions, size=len(questions))
                except Exception as e:
                    message = f"Error processing the TXT file: {e}"
                    return render_template("ack.html", message=message)
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not a text file.")
        
        # Handle YouTube video URL
        elif option == 'video' and request.form.get('video'):
            video_url = request.form['video']
            video_id = video_url.split('v=')[-1]
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                formatter = TextFormatter()
                formatted_transcript = formatter.format_transcript(transcript)
                
                transcript_filename = f"{video_id}_transcript.txt"
                with open(transcript_filename, 'w', encoding="utf-8") as transcript_file:
                    transcript_file.write(formatted_transcript)
                
                return render_template("select_summary.html", filename=transcript_filename)
            except Exception as e:
                message = f"Error fetching transcript for the video: {e}"
                return render_template("ack.html", message=message)
            
        elif option == 'ppt' and 'ppt' in request.files and request.files['ppt'].filename != '':
            ppt_file = request.files['ppt']
            filename = ppt_file.filename
            ppt_file.save(filename)

            if filename.lower().endswith(('pptx')):

                def extract_text_from_ppt(filename):
                    prs = Presentation(filename)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    text.append(paragraph.text)
                    return "\n".join(text)

                try:
                    extracted_text = extract_text_from_ppt(filename)
                    text_filename = f"{filename}_text.txt"
                    print(text_filename, "\n\n\n")
                    with open(text_filename, "w", encoding="utf-8") as text_file:
                        text_file.write(extracted_text)
                    return render_template("select_summary.html", filename=text_filename)
                except Exception as e:
                    message = f"Error processing the PPT file: {e}"
                    return render_template("ack.html", message=message)
                # except Exception as e:
                #     message = f"Error processing the PPT file: {e}"

                
                # print("\n\n\n", filename, "\n\n\n")
                
                # try:
                #     extracted_text = textract.process(filename).decode('utf-8')
                #     text_filename = f"{filename}_text.txt"
                #     print(text_filename, "\n\n\n")
                #     with open(text_filename, "w", encoding="utf-8") as text_file:
                #         text_file.write(extracted_text)
                #     return render_template("select_summary.html", filename=text_filename)
                # except Exception as e:
                #     message = f"Error processing the PPT file: {e}"
                #     return render_template("ack.html", message=message)
            
            else:
                return render_template("ack.html", message="File uploaded successfully but it is not in the form of pptx. Please make sure it is a .pptx file")

        else:
            message = "Please select a valid option and submit the required information."
            return render_template("ack.html", message=message)



@app.route('/summarize', methods=['POST'])
def summarize():
    print("\n\n\n\nHello\n\n\n\n")
    filename = request.form.get('filename')
    summary_type = request.form.get('summary_type')
    summary_length = request.form.get('summary_length')

    print("\n\n\nWagwan\n\n\n")

    with open(filename, 'r', encoding="utf-8") as file:
        text = file.read()

    max_input_length = 512  # Model's input token limit

    if summary_length == 'short':
        max_output_length = 50
    elif summary_length == 'medium':
        max_output_length = 130
    elif summary_length == 'long':
        max_output_length = 250
    else:
        max_output_length = 130

    truncated_text = text[:max_input_length]

    try:
        if summary_type == 'abstractive':
            summary = bart_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
        elif summary_type == 'extractive':
            summary = bert_summarizer(truncated_text, max_length=max_output_length, min_length=max_output_length // 2, do_sample=False)[0]['summary_text']
        else:
            summary = "Invalid summary type selected."

        summary_filename = f"{filename}_summary.txt"
        with open(summary_filename, 'w', encoding="utf-8") as summary_file:
            summary_file.write(summary)

        message = f"Summary generated and saved to {summary_filename} using {summary_type.upper()} with {summary_length} length."
    except Exception as e:
        message = f"Error during summarization: {e}"

    return render_template("ack.html", message=message)



@app.route('/quiz', methods=['GET', 'POST'])
def quiz():

    UPLOAD_STATUS = False
    questions = dict()

    if request.method == 'POST':
        print("\n\n\n\n\n")
        print("Reached Debug Point 2")
        print("\n\n\n\n\n")
        try:
            uploaded_file = request.files['quiz']
            filename = uploaded_file.filename
            uploaded_file.save(filename)

            with open(filename, 'r') as file:
                text = file.read()

            questions = txt2questions(text)

            print("\n\n\n\n\n")
            print(text)
            print("\n\n\n\n\n")

            # File upload + convert success
            if text is not None:
                UPLOAD_STATUS = True
        except Exception as e:
            print("\n\n\n\n\nError caught:", e, "\n\n\n\n\n")
    return render_template(
        'quiz.html',
        uploaded=UPLOAD_STATUS,
        questions=questions,
        size=len(questions)
    )


@app.route('/results', methods=['POST', 'GET'])
def result():
    # Count correct answers
    correct_q = 0
    for q_num, data in questions.items():
        user_answer = request.form.get(f'question{q_num}')
        print("\nUser answer: ", user_answer)
        print("\nActual answer: ", data['answer'],"\n\n")
        if user_answer == data['answer']:
            correct_q += 1
    return render_template('results.html', total=len(questions), correct=correct_q)


if __name__ == '__main__':
    app.run(debug=True, port=5000)
